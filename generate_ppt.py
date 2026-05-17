"""
Chaukas - National Road Safety Hackathon 2026 (IIT Madras)
Professional 7-Slide PPT Generator
Problem Statement: RoadSoS (Emergency Response)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
BG_DARK = RGBColor(0x0A, 0x0A, 0x14)
BG_CARD = RGBColor(0x12, 0x12, 0x1E)
ACCENT_RED = RGBColor(0xFF, 0x2A, 0x2A)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_ORANGE = RGBColor(0xFB, 0x93, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
LIGHT_GRAY = RGBColor(0xD1, 0xD5, 0xDB)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ══════════════════════════════════════════════════════════════
# SLIDE 1: WELCOME / TITLE
# ══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide1)

# Top bar
add_shape(slide1, Inches(0), Inches(0), Inches(13.333), Pt(5), ACCENT_RED)

# Hackathon badge
add_text(slide1, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
         "NATIONAL ROAD SAFETY HACKATHON 2026  •  IIT MADRAS", 14, ACCENT_CYAN, True, PP_ALIGN.CENTER)

# Main Title
add_text(slide1, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
         "CHAUKAS", 72, WHITE, True, PP_ALIGN.CENTER)

add_accent_line(slide1, Inches(5.5), Inches(3.2), Inches(2.3), ACCENT_RED)

add_text(slide1, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
         "India's AI-Powered Road Safety Command System", 28, ACCENT_CYAN, False, PP_ALIGN.CENTER)

# Problem statement badge
badge = add_shape(slide1, Inches(4.2), Inches(4.6), Inches(5), Inches(0.6), BG_CARD, ACCENT_PURPLE)
add_text(slide1, Inches(4.2), Inches(4.65), Inches(5), Inches(0.5),
         "PROBLEM STATEMENT: RoadSoS — Emergency Response", 15, ACCENT_PURPLE, True, PP_ALIGN.CENTER)

# Team info
add_text(slide1, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
         "Team: Outliers  |  Python • React • FastAPI • Leaflet.js", 16, GRAY, False, PP_ALIGN.CENTER)

add_text(slide1, Inches(1), Inches(6.3), Inches(11), Inches(0.4),
         "Offline-First  •  Open Source  •  No Cloud Dependencies", 13, GRAY, False, PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide1, Inches(0), Inches(7.3), Inches(13.333), Pt(5), ACCENT_CYAN)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: PROBLEM & SOLUTION
# ══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)
add_shape(slide2, Inches(0), Inches(0), Inches(13.333), Pt(3), ACCENT_RED)

add_text(slide2, Inches(0.6), Inches(0.3), Inches(5), Inches(0.5),
         "THE PROBLEM", 14, ACCENT_RED, True)
add_text(slide2, Inches(0.6), Inches(0.7), Inches(6), Inches(0.6),
         "Why Road Safety Needs a Command Center", 30, WHITE, True)

# Problem stats cards
problems = [
    ("1.78 Lakh+", "Road deaths in India annually\n(MoRTH 2023)", ACCENT_RED),
    ("~50%", "Deaths due to delayed\nemergency response", ACCENT_ORANGE),
    ("No Single Platform", "For trauma centers, ambulance,\npolice & towing access", ACCENT_PURPLE),
]
for i, (stat, desc, color) in enumerate(problems):
    x = Inches(0.6 + i * 4.1)
    card = add_shape(slide2, x, Inches(1.6), Inches(3.8), Inches(1.6), BG_CARD, color)
    add_text(slide2, x + Inches(0.3), Inches(1.75), Inches(3.2), Inches(0.6), stat, 32, color, True)
    add_text(slide2, x + Inches(0.3), Inches(2.4), Inches(3.2), Inches(0.7), desc, 13, LIGHT_GRAY)

# Solution section
add_accent_line(slide2, Inches(0.6), Inches(3.6), Inches(12), ACCENT_CYAN)
add_text(slide2, Inches(0.6), Inches(3.9), Inches(5), Inches(0.5),
         "OUR SOLUTION: CHAUKAS", 14, ACCENT_CYAN, True)

add_text(slide2, Inches(0.6), Inches(4.4), Inches(12), Inches(0.8),
         "A unified, offline-first road safety command platform that provides location-based access to trauma centers, ambulances, police, and towing services — with AI-powered severity assessment and ML-driven accident hotspot detection.",
         16, LIGHT_GRAY)

# Solution pillars
pillars = [
    ("🧠 AI Severity Engine", "Auto-assesses risk using\nGPS + weather + incidents", ACCENT_CYAN),
    ("🗺️ ML Hotspot Detection", "Identifies accident blackspots\nnear user location", ACCENT_GREEN),
    ("📡 Offline-First Design", "Works without internet.\nLocal DB, no cloud needed.", ACCENT_ORANGE),
    ("⚡ Escalation Engine", "NORMAL → WATCH →\nPREPAREDNESS → CRISIS", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(pillars):
    x = Inches(0.6 + i * 3.15)
    card = add_shape(slide2, x, Inches(5.4), Inches(2.9), Inches(1.7), BG_CARD, color)
    add_text(slide2, x + Inches(0.2), Inches(5.55), Inches(2.5), Inches(0.5), title, 14, color, True)
    add_text(slide2, x + Inches(0.2), Inches(6.05), Inches(2.5), Inches(0.7), desc, 12, GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: TECHNICAL ARCHITECTURE
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)
add_shape(slide3, Inches(0), Inches(0), Inches(13.333), Pt(3), ACCENT_CYAN)

add_text(slide3, Inches(0.6), Inches(0.3), Inches(5), Inches(0.5),
         "TECHNICAL ARCHITECTURE", 14, ACCENT_CYAN, True)
add_text(slide3, Inches(0.6), Inches(0.7), Inches(8), Inches(0.6),
         "System Design & Technology Stack", 30, WHITE, True)

# Frontend box
add_shape(slide3, Inches(0.5), Inches(1.6), Inches(3.8), Inches(5.3), BG_CARD, ACCENT_CYAN)
add_text(slide3, Inches(0.7), Inches(1.7), Inches(3.4), Inches(0.4), "FRONTEND", 14, ACCENT_CYAN, True)
fe_items = [
    "React 18 + Vite 7",
    "Leaflet.js (Maps & Heatmaps)",
    "Framer Motion (Animations)",
    "IndexedDB (Offline Storage)",
    "PWA + Service Worker",
    "3D Globe Visualization",
    "Real-time WebSocket feeds",
    "Responsive Mobile-First UI",
]
for j, item in enumerate(fe_items):
    add_text(slide3, Inches(0.7), Inches(2.2 + j*0.55), Inches(3.4), Inches(0.4),
             f"▸ {item}", 12, LIGHT_GRAY)

# Backend box
add_shape(slide3, Inches(4.7), Inches(1.6), Inches(3.8), Inches(5.3), BG_CARD, ACCENT_PURPLE)
add_text(slide3, Inches(4.9), Inches(1.7), Inches(3.4), Inches(0.4), "BACKEND (Python)", 14, ACCENT_PURPLE, True)
be_items = [
    "FastAPI + Uvicorn",
    "ML Hotspot Engine (DBSCAN)",
    "AI Severity Scoring Engine",
    "Escalation State Machine",
    "Gemini Vision (Photo AI)",
    "News Intelligence Scraper",
    "SQLite (Local Persistence)",
    "Mock Supabase Client",
]
for j, item in enumerate(be_items):
    add_text(slide3, Inches(4.9), Inches(2.2 + j*0.55), Inches(3.4), Inches(0.4),
             f"▸ {item}", 12, LIGHT_GRAY)

# Data flow box
add_shape(slide3, Inches(8.9), Inches(1.6), Inches(4), Inches(5.3), BG_CARD, ACCENT_GREEN)
add_text(slide3, Inches(9.1), Inches(1.7), Inches(3.6), Inches(0.4), "DATA FLOW", 14, ACCENT_GREEN, True)
flow_items = [
    "1. User GPS → Browser API",
    "2. Location → Severity Engine",
    "3. Weather + Incidents → Risk",
    "4. Risk Score → Escalation",
    "5. Incidents → ML Clustering",
    "6. Clusters → Heatmap Grid",
    "7. News → AI Categorization",
    "8. All → Admin Dashboard",
]
for j, item in enumerate(flow_items):
    add_text(slide3, Inches(9.1), Inches(2.2 + j*0.55), Inches(3.6), Inches(0.4),
             f"{item}", 12, LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: KEY FEATURES (RoadSoS Alignment)
# ══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)
add_shape(slide4, Inches(0), Inches(0), Inches(13.333), Pt(3), ACCENT_GREEN)

add_text(slide4, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
         "KEY FEATURES — ALIGNED TO RoadSoS CRITERIA", 14, ACCENT_GREEN, True)
add_text(slide4, Inches(0.6), Inches(0.7), Inches(10), Inches(0.6),
         "Location-Based Emergency Response Platform", 30, WHITE, True)

features = [
    ("🚨 Crisis Reporting", "Report accidents with photo, location, severity.\nAI auto-classifies incident type via Gemini Vision.\nGenerates unique incident ID for tracking.",
     ACCENT_RED, "Reliability ✓"),
    ("🏥 Emergency Contacts", "One-tap access to nearest trauma centers,\nambulances, police, and towing services.\nLocation-aware, prioritized by distance.",
     ACCENT_CYAN, "Contacts Fetched ✓"),
    ("🧠 AI Severity Engine", "Auto-assesses real-time risk at your GPS location.\nFactors: weather, time, incident density, road type.\nOutputs: LOW → ELEVATED → HIGH → CRITICAL.",
     ACCENT_PURPLE, "Innovation ✓"),
    ("🗺️ ML Hotspot Detection", "DBSCAN spatial clustering on accident data.\nGenerates color-coded heatmap grid zones.\nSafe / Caution / Warning / Danger classification.",
     ACCENT_GREEN, "Innovation ✓"),
    ("⚡ Escalation Engine", "4-state machine: NORMAL → WATCH →\nPREPAREDNESS → CRISIS. Auto-triggers from\nSeverity AI. Real-time 30s refresh cycle.",
     ACCENT_ORANGE, "Service Integration ✓"),
    ("📡 Offline-First", "Full functionality without internet.\nLocal JSON DB + SQLite + IndexedDB.\nPWA with service worker caching.",
     RGBColor(0xA7, 0x8B, 0xFA), "Offline Robustness ✓"),
]

for i, (title, desc, color, criteria) in enumerate(features):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.5 + row * 2.9)
    card = add_shape(slide4, x, y, Inches(3.9), Inches(2.5), BG_CARD, color)
    add_text(slide4, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.4), title, 15, color, True)
    add_text(slide4, x + Inches(0.2), y + Inches(0.6), Inches(3.5), Inches(1.2), desc, 11, GRAY)
    # Criteria badge
    badge = add_shape(slide4, x + Inches(0.2), y + Inches(1.95), Inches(2.2), Inches(0.35), RGBColor(0x1A, 0x1A, 0x2E), color)
    add_text(slide4, x + Inches(0.2), y + Inches(1.97), Inches(2.2), Inches(0.3), criteria, 10, color, True, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: DEMO / SCREENSHOTS LAYOUT
# ══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)
add_shape(slide5, Inches(0), Inches(0), Inches(13.333), Pt(3), ACCENT_ORANGE)

add_text(slide5, Inches(0.6), Inches(0.3), Inches(5), Inches(0.5),
         "LIVE DEMO WALKTHROUGH", 14, ACCENT_ORANGE, True)
add_text(slide5, Inches(0.6), Inches(0.7), Inches(10), Inches(0.6),
         "User Journey — From Incident to Resolution", 30, WHITE, True)

steps = [
    ("STEP 1", "Landing & Login", "User arrives at Chaukas.\nSecure authentication via\nlocal credential system.\nGPS auto-acquired.", ACCENT_CYAN),
    ("STEP 2", "Report Incident", "Upload photo/video of accident.\nGemini AI analyzes severity.\nAuto-fills location from GPS.\nUnique incident ID generated.", ACCENT_RED),
    ("STEP 3", "Severity Analysis", "AI engine processes:\n• Weather conditions\n• Nearby incident density\n• Time-of-day risk factors\nOutputs: risk band + confidence.", ACCENT_PURPLE),
    ("STEP 4", "Escalation Response", "State machine transitions:\nNORMAL → WATCH →\nPREPAREDNESS → CRISIS\nRecommended actions shown.", ACCENT_ORANGE),
    ("STEP 5", "Hotspot Map", "ML clustering reveals danger\nzones near user location.\nColor-coded heatmap grid.\n25/50/100 km radius options.", ACCENT_GREEN),
    ("STEP 6", "Admin Command", "Dashboard with live stats,\nincident management, user\ntracking, system analytics.\nFull command center view.", RGBColor(0xA7, 0x8B, 0xFA)),
]

for i, (step, title, desc, color) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.5 + row * 2.9)
    card = add_shape(slide5, x, y, Inches(3.9), Inches(2.5), BG_CARD, color)
    # Step number
    add_text(slide5, x + Inches(0.2), y + Inches(0.1), Inches(1), Inches(0.3), step, 10, color, True)
    add_text(slide5, x + Inches(0.2), y + Inches(0.4), Inches(3.5), Inches(0.4), title, 16, WHITE, True)
    add_accent_line(slide5, x + Inches(0.2), y + Inches(0.85), Inches(1.5), color)
    add_text(slide5, x + Inches(0.2), y + Inches(1.0), Inches(3.5), Inches(1.3), desc, 11, GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: IMPACT & INNOVATION
# ══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6)
add_shape(slide6, Inches(0), Inches(0), Inches(13.333), Pt(3), ACCENT_PURPLE)

add_text(slide6, Inches(0.6), Inches(0.3), Inches(5), Inches(0.5),
         "IMPACT & INNOVATION", 14, ACCENT_PURPLE, True)
add_text(slide6, Inches(0.6), Inches(0.7), Inches(10), Inches(0.6),
         "Why Chaukas Stands Out", 30, WHITE, True)

# Innovation highlights
innovations = [
    ("Zero Cloud Dependency", "Entire platform runs locally. No AWS, no Firebase,\nno Supabase needed. Perfect for remote highway\ncommand centers with poor connectivity.", ACCENT_GREEN),
    ("AI-Driven, No Manual Input", "Severity assessment is fully autonomous —\nuses GPS + weather + incident correlation.\nZero manual data entry required from users.", ACCENT_CYAN),
    ("Real-Time ML Hotspots", "DBSCAN clustering identifies accident blackspots\nin real-time. Dynamic heatmap updates every 30s.\nHelps traffic police deploy resources proactively.", ACCENT_ORANGE),
    ("Multi-Feature Integration", "9+ interconnected features in one platform:\nSeverity → Escalation → Hotspot → News →\nAdmin — all sharing location context.", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(innovations):
    x = Inches(0.5 + (i % 2) * 6.3)
    y = Inches(1.5 + (i // 2) * 2.2)
    card = add_shape(slide6, x, y, Inches(6), Inches(1.9), BG_CARD, color)
    add_text(slide6, x + Inches(0.3), y + Inches(0.15), Inches(5.4), Inches(0.4), title, 16, color, True)
    add_text(slide6, x + Inches(0.3), y + Inches(0.6), Inches(5.4), Inches(1.1), desc, 13, LIGHT_GRAY)

# Impact metrics
add_accent_line(slide6, Inches(0.6), Inches(5.9), Inches(12), ACCENT_PURPLE)
add_text(slide6, Inches(0.6), Inches(6.1), Inches(3), Inches(0.3), "PROJECTED IMPACT", 12, ACCENT_PURPLE, True)

metrics = [
    ("30% Faster", "Emergency\nResponse Time", ACCENT_GREEN),
    ("100% Offline", "Operational in\nNo-Network Zones", ACCENT_CYAN),
    ("9+ Features", "Integrated in\nOne Platform", ACCENT_ORANGE),
    ("Real-Time", "ML-Powered\nRisk Assessment", ACCENT_PURPLE),
]
for i, (val, label, color) in enumerate(metrics):
    x = Inches(0.6 + i * 3.15)
    add_text(slide6, x, Inches(6.35), Inches(2.8), Inches(0.4), val, 22, color, True)
    add_text(slide6, x, Inches(6.75), Inches(2.8), Inches(0.5), label, 11, GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: THANK YOU
# ══════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7)
add_shape(slide7, Inches(0), Inches(0), Inches(13.333), Pt(5), ACCENT_CYAN)

add_text(slide7, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "THANK YOU", 64, WHITE, True, PP_ALIGN.CENTER)

add_accent_line(slide7, Inches(5.5), Inches(2.8), Inches(2.3), ACCENT_RED)

add_text(slide7, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
         "Chaukas — Making Indian Roads Safer with AI", 24, ACCENT_CYAN, False, PP_ALIGN.CENTER)

# Tech summary
add_text(slide7, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
         "Python  •  FastAPI  •  React  •  Leaflet.js  •  Gemini AI  •  SQLite", 16, GRAY, False, PP_ALIGN.CENTER)

# Contact card
card = add_shape(slide7, Inches(3.5), Inches(5.0), Inches(6.3), Inches(1.5), BG_CARD, ACCENT_CYAN)
add_text(slide7, Inches(3.5), Inches(5.15), Inches(6.3), Inches(0.4),
         "Team Outliers", 20, WHITE, True, PP_ALIGN.CENTER)
add_text(slide7, Inches(3.5), Inches(5.55), Inches(6.3), Inches(0.3),
         "National Road Safety Hackathon 2026 — IIT Madras", 13, ACCENT_CYAN, False, PP_ALIGN.CENTER)
add_text(slide7, Inches(3.5), Inches(5.9), Inches(6.3), Inches(0.3),
         "Problem Statement: RoadSoS (Emergency Response)", 12, GRAY, False, PP_ALIGN.CENTER)
add_text(slide7, Inches(3.5), Inches(6.2), Inches(6.3), Inches(0.3),
         "100% Open Source  •  100% Offline  •  100% Local", 12, ACCENT_GREEN, True, PP_ALIGN.CENTER)

add_shape(slide7, Inches(0), Inches(7.3), Inches(13.333), Pt(5), ACCENT_RED)

# Save
output_path = os.path.join(os.path.dirname(__file__), "Chaukas_RoadSafety_Hackathon_2026.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
